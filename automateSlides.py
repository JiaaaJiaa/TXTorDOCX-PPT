from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN  # Import PP_ALIGN
import os

# Function to read content from a .txt file
def read_txt(file_path):
    with open(file_path, 'r') as file:
        return file.readlines()

# Function to read content from a .docx file
def read_docx(file_path):
    from docx import Document
    document = Document(file_path)
    return [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]

# Function to create slides with 2 lines per slide
def create_ppt_from_file(file_path, output_file):
    # Read file content
    if file_path.endswith('.txt'):
        lines = read_txt(file_path)
    elif file_path.endswith('.docx'):
        lines = read_docx(file_path)
    else:
        raise ValueError("Unsupported file type. Use .txt or .docx files.")
    
    # Initialize PowerPoint
    presentation = Presentation()
    
    # Set slide dimensions to 1920x1080
    presentation.slide_width = Inches(20)
    presentation.slide_height = Inches(11.25)
    
    # Split lines into groups of 2 per slide
    for i in range(0, len(lines), 2):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])  # Blank layout
        # Position the text box in the lower third area
        text_box = slide.shapes.add_textbox(Inches(1), Inches(7), Inches(18), Inches(3))  # Textbox position
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        # Add two lines to the slide
        for j in range(2):
            if i + j < len(lines):  # Ensure no IndexError
                if j > 0:
                    p = text_frame.add_paragraph()
                else:
                    p = text_frame.paragraphs[0]
                
                p.text = lines[i + j]
                # Set font properties
                p.font.size = Pt(44)  # Set font size to 44
                p.font.bold = True  # Make the text bold
                p.font.color.rgb = RGBColor(255, 255, 255)  # Set the text color to white
                p.alignment = PP_ALIGN.CENTER  # Center align the text
    
    # Save the presentation
    presentation.save(output_file)
    print(f"Presentation saved as '{output_file}'")

# Example Usage
input_file = "sample.docx"  # Replace with your .txt or .docx file path
output_pptx = "output_presentation.pptx"
create_ppt_from_file(input_file, output_pptx)